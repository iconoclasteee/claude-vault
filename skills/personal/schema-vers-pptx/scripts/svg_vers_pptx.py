# -*- coding: utf-8 -*-
"""Convertit un SVG en PowerPoint à formes natives — et répare ce que svg2pptx casse.

svg2pptx recrée bien la géométrie et les textes simples. Trois défauts, aucun ne
lève d'erreur, tous visibles seulement au rendu :

  1. Les <tspan> sont ignorés. Une ligne faite de segments de couleurs différentes
     ressort en un seul bloc SANS couleur, donc noir. Sur un fond sombre, le texte
     blanc devient noir sur noir. Réparation : relire le SVG, retrouver chaque ligne
     à son abscisse, reconstruire les runs.
  2. Les corps de police sont 4/3 trop grands. Les formes sont posées à 0,632 pt/unité
     (largeur du SVG ramenée à la largeur de la diapo), le texte à 0,842 pt/unité.
     Résultat : textes débordant de leur bloc, chevauchements, retours à la ligne.
  3. La pile de polices CSS ("Inter, 'Segoe UI', …") est passée telle quelle comme
     nom de police. PowerPoint ne la trouve pas et retombe sur son thème.

Usage :
    ~/.venvs/svg2pptx/bin/python svg_vers_pptx.py mon-schema.svg          # -> mon-schema.pptx
    ~/.venvs/svg2pptx/bin/python svg_vers_pptx.py mon-schema.svg sortie.pptx
    ~/.venvs/svg2pptx/bin/python svg_vers_pptx.py mon-schema.svg --pdf    # + PDF de contrôle

Prérequis : un venv, PAS le Python système — sur macOS avec Homebrew, `pip install`
échoue (PEP 668, externally-managed-environment) :
    python3 -m venv ~/.venvs/svg2pptx
    ~/.venvs/svg2pptx/bin/pip install svg2pptx python-pptx lxml
Le PDF de contrôle exige LibreOffice (`soffice` dans le PATH). Il s'écrit à côté du
PPTX produit, sans option de redirection.

⚠️ Regarder le PDF. Le script peut annoncer « 47/47 lignes recolorées » sur un rendu
   cassé : il ne sait pas ce que le schéma est censé montrer.
"""
import argparse
import os
import subprocess
import sys

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu
from svg2pptx import svg_to_pptx

POLICE, POLICE_MONO = "Arial", "Consolas"
CORRECTION_CORPS = 0.75          # 0,632 / 0,842 — voir défaut n°2
SVG_NS = "{http://www.w3.org/2000/svg}"


def segments_du_svg(chemin):
    """Relève les lignes à tspans : {(texte complet, x arrondi) -> [(bout, couleur, gras)]}.

    Renvoie aussi les clés ambiguës : deux lignes de même texte à la même abscisse mais
    de couleurs différentes (même colonne, lignes différentes) sont indiscernables ici,
    et la réparation appliquerait les mauvaises couleurs à l'une des deux. Rare, mais
    silencieux — donc signalé.
    """
    racine = etree.parse(chemin).getroot()
    releve, ambigues = {}, set()
    for t in racine.iter(SVG_NS + "text"):
        bouts = t.findall(SVG_NS + "tspan")
        if not bouts:
            continue
        segs = [(b.text or "", b.get("fill"), b.get("font-weight") in ("600", "700"))
                for b in bouts]
        cle = ("".join(s for s, _, _ in segs), round(float(t.get("x"))))
        if cle in releve and releve[cle] != segs:
            ambigues.add(cle)
        releve[cle] = segs
    return releve, ambigues


def largeur_svg(chemin):
    return float(etree.parse(chemin).getroot().get("width"))


def repare(pptx_path, releve, px_par_emu, centrer=True):
    prs = Presentation(pptx_path)
    diapo = prs.slides[0]
    rendus = debordements = 0

    for forme in diapo.shapes:
        if not forme.has_text_frame:
            continue
        cadre = forme.text_frame
        cadre.word_wrap = False

        for para in cadre.paragraphs:                       # défauts 2 et 3
            for run in para.runs:
                nom = run.font.name or ""
                run.font.name = POLICE_MONO if ("Mono" in nom or "Menlo" in nom) else POLICE
                if run.font.size:
                    run.font.size = Emu(int(run.font.size * CORRECTION_CORPS))

        segs = releve.get((cadre.text, round(forme.left * px_par_emu)))   # défaut 1
        if segs and cadre.paragraphs:
            para = cadre.paragraphs[0]
            taille = para.runs[0].font.size if para.runs else None
            for run in list(para.runs):
                run._r.getparent().remove(run._r)
            for bout, couleur, gras in segs:
                run = para.add_run()
                run.text = bout
                run.font.size = taille
                run.font.name = POLICE
                run.font.bold = gras
                if couleur:
                    run.font.color.rgb = RGBColor.from_string(couleur.lstrip("#").upper())
            rendus += 1

        if forme.left is not None and forme.width and forme.left + forme.width > prs.slide_width:
            forme.width = prs.slide_width - forme.left
            debordements += 1

    if centrer:
        # Un SVG large posé sur un 16:9 laisse une bande blanche en bas : on répartit.
        bas = max((f.top + f.height) for f in diapo.shapes
                  if f.top is not None and f.height is not None)
        marge = (prs.slide_height - bas) // 2
        if marge > 0:
            for forme in diapo.shapes:
                if forme.top is not None:
                    forme.top = forme.top + marge

    prs.save(pptx_path)
    return rendus, debordements


def controle_pdf(pptx_path):
    """Exporte un PDF pour vérification visuelle. À REGARDER, pas seulement à produire."""
    dossier = os.path.dirname(os.path.abspath(pptx_path))
    subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", dossier,
                    pptx_path], check=True, capture_output=True)
    return os.path.splitext(pptx_path)[0] + ".pdf"


def convertit(svg, pptx=None, pdf=False, centrer=True):
    pptx = pptx or os.path.splitext(svg)[0] + ".pptx"
    svg_to_pptx(svg, pptx)
    releve, ambigues = segments_du_svg(svg)
    px_par_emu = largeur_svg(svg) / Presentation(pptx).slide_width
    rendus, debordements = repare(pptx, releve, px_par_emu, centrer)
    print(f"{os.path.basename(pptx)} | {rendus}/{len(releve)} lignes recolorées, "
          f"{debordements} débordement(s) rentré(s)")
    if rendus < len(releve):
        print(f"  ⚠ {len(releve) - rendus} ligne(s) du SVG non retrouvée(s) dans le PPTX : "
              f"elles resteront en noir. À examiner au rendu.")
    for texte, x in sorted(ambigues):
        print(f"  ⚠ ligne ambiguë (même texte, même x={x}) : « {texte[:50]} » — "
              f"couleurs possiblement inversées entre deux occurrences.")
    if pdf:
        print("contrôle visuel :", controle_pdf(pptx))
    return pptx


if __name__ == "__main__":
    p = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    p.add_argument("svg")
    p.add_argument("pptx", nargs="?")
    p.add_argument("--pdf", action="store_true", help="exporte un PDF de contrôle (LibreOffice)")
    p.add_argument("--sans-centrage", action="store_true", help="laisse le schéma en haut")
    a = p.parse_args()
    if not os.path.exists(a.svg):
        sys.exit(f"introuvable : {a.svg}")
    convertit(a.svg, a.pptx, a.pdf, not a.sans_centrage)
