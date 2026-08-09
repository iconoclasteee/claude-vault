# -*- coding: utf-8 -*-
"""Plomberie pour dessiner un schéma directement en PowerPoint, en formes natives.

Le problème que ce module résout : python-pptx raisonne en EMU et en boîtes de texte,
alors qu'un schéma se pense en coordonnées abstraites et en lignes de base. `Toile`
fait la traduction, et regroupe les formes pour que le résultat reste manipulable
à la souris.

    from pptx_kit import Toile

    t = Toile()                                  # 1520 unités de large, diapo 16:9
    g = t.groupe()                               # un groupe = un objet déplaçable
    t.bloc(g, 32, 100, 264, 200, "EEF4FA", "A9C5DF")
    t.texte(g, 50, 130, "ÉTAGE 1", 10, "7D9CBB", gras=True, spc=1.4)
    t.texte(g, 50, 160, [("30 critères en ", 0), ("8 piliers", 1)], 12.4,
            {False: "22303F", True: "1B5C96"})
    t.enregistre("schema.pptx")

Toutes les couleurs sont des hexadécimaux SANS dièse ("EEF4FA").
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


class Toile:
    """Une diapositive dessinée dans un repère abstrait.

    `unites` fixe la largeur du repère de travail ; la hauteur en découle du ratio.
    Travailler à 1520 unités permet de reprendre tel quel un tracé pensé en pixels SVG.
    """

    def __init__(self, unites=1520, ratio=(16, 9), police="Arial", police_mono="Consolas"):
        self.unites = unites
        self.hauteur_unites = unites * ratio[1] / ratio[0]
        self.police, self.police_mono = police, police_mono
        self.prs = Presentation()
        self.prs.slide_width = 12192000
        self.prs.slide_height = int(12192000 * ratio[1] / ratio[0])
        self.diapo = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.u = self.prs.slide_width / unites            # EMU par unité
        self.pt = self.u / 12700                          # points par unité

    # --- conversions ---------------------------------------------------------
    def e(self, v):
        return Emu(int(round(v * self.u)))

    # --- formes --------------------------------------------------------------
    def groupe(self):
        """Un conteneur déplaçable d'un clic. Y ranger tout ce qui fait bloc."""
        return self.diapo.shapes.add_group_shape()

    def bloc(self, hote, x, y, w, h, fond, ligne=None, rayon=7):
        f = hote.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  self.e(x), self.e(y), self.e(w), self.e(h))
        f.adjustments[0] = rayon / min(w, h)
        f.fill.solid()
        f.fill.fore_color.rgb = RGBColor.from_string(fond)
        if ligne:
            f.line.color.rgb = RGBColor.from_string(ligne)
            f.line.width = Pt(0.75)
        else:
            f.line.fill.background()
        f.shadow.inherit = False
        # shadow.inherit ne suffit pas : la forme garde un <a:effectRef> vers l'effet
        # du thème, que certains lecteurs (LibreOffice) appliquent quand même. On
        # neutralise la référence elle-même, sinon les blocs sortent avec une ombre portée.
        style = f._element.find(qn("p:style"))
        if style is not None:
            effet = style.find(qn("a:effectRef"))
            if effet is not None:
                effet.set("idx", "0")
        f.text_frame.text = ""
        return f

    def filet(self, hote, x1, y, x2, couleur, epaisseur=0.75):
        c = hote.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      self.e(x1), self.e(y), self.e(x2), self.e(y))
        c.line.color.rgb = RGBColor.from_string(couleur)
        c.line.width = Pt(epaisseur)
        return c

    def fleche(self, hote, x, y, longueur=17, couleur="7D9CBB"):
        c = hote.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      self.e(x), self.e(y), self.e(x + longueur), self.e(y))
        c.line.color.rgb = RGBColor.from_string(couleur)
        c.line.width = Pt(1.2)
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
        return c

    # --- texte ---------------------------------------------------------------
    def texte(self, hote, x, ligne_base, segments, taille, couleur="22303F", gras=False,
              cadrage="l", police=None, spc=None, largeur=320):
        """Une ligne de texte, posée sur sa ligne de base.

        `segments` : une chaîne, ou une liste de (texte, gras) pour une ligne à
        plusieurs styles. `couleur` : une chaîne, ou {False: normal, True: accentué}.
        `cadrage` : "l" (x = bord gauche), "c" (x = centre), "r" (x = bord droit).
        `spc` : interlettrage en points — ce que la conversion depuis SVG perd.

        Pourquoi une boîte centrée plutôt qu'un calage sur la ligne de base :
        PowerPoint positionne une boîte, pas une ligne de base, et chaque logiciel
        calcule l'ascendante de la police à sa façon. Centrer la boîte sur la hauteur
        d'œil rend le rendu stable d'un lecteur à l'autre.
        """
        pts = taille * self.pt
        h = taille * 1.7
        cy = ligne_base - taille * 0.32
        if cadrage == "l":
            gx, gw, al = x, min(largeur, self.unites - x), PP_ALIGN.LEFT
        elif cadrage == "r":
            gw = min(largeur, x)
            gx, al = x - gw, PP_ALIGN.RIGHT
        else:
            gw = min(largeur, 2 * x, 2 * (self.unites - x))
            gx, al = x - gw / 2, PP_ALIGN.CENTER
        zone = hote.shapes.add_textbox(self.e(gx), self.e(cy - h / 2), self.e(gw), self.e(h))
        cadre = zone.text_frame
        cadre.word_wrap = False         # une ligne de source = une ligne de rendu
        cadre.margin_left = cadre.margin_right = 0
        cadre.margin_top = cadre.margin_bottom = 0
        cadre.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cadre.paragraphs[0]
        para.alignment = al
        for texte, fort in ([(segments, gras)] if isinstance(segments, str) else segments):
            r = para.add_run()
            r.text = texte
            r.font.size = Pt(round(pts, 1))
            r.font.name = police or self.police
            r.font.bold = bool(fort)
            r.font.color.rgb = RGBColor.from_string(
                couleur if isinstance(couleur, str) else couleur[bool(fort)])
            if spc:
                r.font._rPr.set("spc", str(int(spc * 100)))
        return zone

    def pastille(self, hote, x_droite, y_centre, libelle, largeur=112, hauteur=18,
                 fond="12406F", encre="FFFFFF"):
        """Étiquette en gélule, posée à cheval sur le bord d'un bloc."""
        p = self.bloc(hote, x_droite - largeur, y_centre - hauteur / 2, largeur, hauteur,
                      fond, encre, hauteur / 2)
        p.line.width = Pt(1.5)
        cadre = p.text_frame
        cadre.margin_left = cadre.margin_right = 0
        cadre.margin_top = cadre.margin_bottom = 0
        cadre.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cadre.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = libelle
        r.font.size = Pt(round(hauteur * 0.52 * self.pt, 1))
        r.font.name = self.police
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string(encre)
        return p

    # --- sortie --------------------------------------------------------------
    def enregistre(self, chemin):
        self.prs.save(chemin)
        haut = list(self.diapo.shapes)
        groupes = [s for s in haut if s.shape_type == 6]
        total = sum(len(g.shapes) for g in groupes) + (len(haut) - len(groupes))
        return {"objets": len(haut), "groupes": len(groupes), "formes": total}
